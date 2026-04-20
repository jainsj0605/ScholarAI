from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()
    
    # helper for slide background
    def set_slide_background(slide, color_rgb):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*color_rgb)

    # Helper for adding image placeholders
    def add_image_placeholder(slide, top_inch, left_inch, width_inch=6):
        box = slide.shapes.add_shape(
            1, # Rectangle
            Inches(left_inch), Inches(top_inch), Inches(width_inch), Inches(3)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(42, 51, 71) # Dark Slate
        line = box.line
        line.color.rgb = RGBColor(76, 175, 130) # Emerald Green
        line.width = Pt(2)
        
        tf = box.text_frame
        tf.text = "[ INSERT SCREENSHOT HERE ]"
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(255, 255, 255)

    # 1. Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_slide_background(slide, (15, 17, 23)) # Dark Theme
    title = slide.shapes.title
    title.text = "Research Paper Analysis Using Multimodal RAG and Agentic AI"
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title.text_frame.paragraphs[0].font.bold = True

    subtitle = slide.placeholders[1]
    subtitle.text = "by\nAgastya Kundu (23UEC508) | Ayush Negi (23UCC528) | Siddharth Jain (23UEC625)\n\nUnder Guidance of: Dr. Vaibhav Kumar Gupta\nDepartment of ECE, LNMIIT Jaipur"
    subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(180, 180, 180)

    # 2. Motivation
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_background(slide, (15, 17, 23))
    slide.shapes.title.text = "Problem Statement & Motivation"
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(76, 175, 130)
    
    tf = slide.placeholders[1].text_frame
    tf.text = "• Literature Overload: 5,000+ new papers daily in Engineering."
    tf.add_paragraph().text = "• Loss of Context: Figures/Tables are disconnected from text in standard RAG."
    tf.add_paragraph().text = "• The 'Not Reported' Trap: Generic AI misses specific technical results."
    tf.add_paragraph().text = "• Need: An automated technical auditor that understands complex academic PDF layouts."
    for p in tf.paragraphs: p.font.color.rgb = RGBColor(255, 255, 255)

    # 3. Core Architecture
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_background(slide, (15, 17, 23))
    slide.shapes.title.text = "System Architecture: LangGraph"
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(76, 175, 130)
    tf = slide.placeholders[1].text_frame
    tf.text = "• Agentic Orchestration: Moves beyond linear pipelines using LangGraph."
    tf.add_paragraph().text = "• State Machine: Cyclic reasoning between Summarizer, Critic, and Writer nodes."
    for p in tf.paragraphs: p.font.color.rgb = RGBColor(255, 255, 255)
    add_image_placeholder(slide, 3.5, 2)

    # 4. Data Ingestion & Figures
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_background(slide, (15, 17, 23))
    slide.shapes.title.text = "Multimodal Figure Discovery"
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(76, 175, 130)
    tf = slide.placeholders[1].text_frame
    tf.text = "• Spatial Grounding: Extracts text directly above and below figures."
    tf.add_paragraph().text = "• Multi-Engine Parsing: Combining PyMuPDF and fitz for vector fidelity."
    tf.add_paragraph().text = "• Vision Grounding: Sends image + spatial context to AI to ensure zero-hallucination."
    for p in tf.paragraphs: p.font.color.rgb = RGBColor(255, 255, 255)
    add_image_placeholder(slide, 4.0, 2)

    # 5. Retrieval Strategy (RAG)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_background(slide, (15, 17, 23))
    slide.shapes.title.text = "Engineering-Grade RAG"
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(76, 175, 130)
    tf = slide.placeholders[1].text_frame
    tf.text = "• Multi-Query Expansion: Splits search into Problem, Method, and Results."
    tf.add_paragraph().text = "• Retrieval Depth: k=5 per query for 360-degree technical coverage."
    tf.add_paragraph().text = "• Sliding Window: 500-char chunks for precise semantic matching."
    for p in tf.paragraphs: p.font.color.rgb = RGBColor(255, 255, 255)

    # 6. Technical Summarization
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_background(slide, (15, 17, 23))
    slide.shapes.title.text = "Automated Technical Audit"
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(76, 175, 130)
    tf = slide.placeholders[1].text_frame
    tf.text = "• Persona: AI acts as a Senior Technical Reviewer."
    tf.add_paragraph().text = "• LaTeX Support: Native mathematical equation rendering ($)."
    tf.add_paragraph().text = "• Structured Sections: Architecture, Performance, Simulation, and Results."
    for p in tf.paragraphs: p.font.color.rgb = RGBColor(255, 255, 255)
    add_image_placeholder(slide, 4.0, 2)

    # 7. Comparison Matrix
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_background(slide, (15, 17, 23))
    slide.shapes.title.text = "Hybrid Cross-Engine Discovery"
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(76, 175, 130)
    tf = slide.placeholders[1].text_frame
    tf.text = "• Parallel Search: ArXiv, OpenAlex, Semantic Scholar, and CrossRef."
    tf.add_paragraph().text = "• Filtering: Prioritizes Engineering venues (IEEE, Springer, Elsevier)."
    tf.add_paragraph().text = "• Comparative Matrix: Side-by-side gap analysis of 6 related works."
    for p in tf.paragraphs: p.font.color.rgb = RGBColor(255, 255, 255)
    add_image_placeholder(slide, 4.0, 2)

    # 8. Improvement Pipeline
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_background(slide, (15, 17, 23))
    slide.shapes.title.text = "Technical Improvement Pipeline"
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(76, 175, 130)
    tf = slide.placeholders[1].text_frame
    tf.text = "• Detection: Identified 5 critical technical gaps in current research."
    tf.add_paragraph().text = "• Rewriting: AI-driven 'Pull Request' to Sound like the original author."
    tf.add_paragraph().text = "• Style Matching: Preserves nuances while strengthening methodology."
    for p in tf.paragraphs: p.font.color.rgb = RGBColor(255, 255, 255)
    add_image_placeholder(slide, 4.0, 2)

    # 9. Technical Stack
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_background(slide, (15, 17, 23))
    slide.shapes.title.text = "Inference & Data Infrastructure"
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(76, 175, 130)
    tf = slide.placeholders[1].text_frame
    tf.text = "• Inference: Groq LPU (Language Processing Unit) for ultra-fast response."
    tf.add_paragraph().text = "• Vector DB: FAISS for high-performance retrieval."
    tf.add_paragraph().text = "• Frontend: Streamlit for interactive dashboard execution."
    for p in tf.paragraphs: p.font.color.rgb = RGBColor(255, 255, 255)

    # 10. Reliability System
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_background(slide, (15, 17, 23))
    slide.shapes.title.text = "Self-Healing Reliability: 3-Tier Fallback"
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(76, 175, 130)
    tf = slide.placeholders[1].text_frame
    tf.text = "• Layer 1: GPT-OSS-120B (High Reasoning)."
    tf.add_paragraph().text = "• Layer 2: Llama-3.3-70B (Reliable backup for summaries)."
    tf.add_paragraph().text = "• Layer 3: Llama-3.1-8B (High-TPM speed tier for Q&A)."
    tf.add_paragraph().text = "• Result: Zero-downtime demo handling 100k+ daily token limits."
    for p in tf.paragraphs: p.font.color.rgb = RGBColor(255, 255, 255)

    # 11. Results & Conclusion
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_background(slide, (15, 17, 23))
    slide.shapes.title.text = "Conclusion"
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(76, 175, 130)
    tf = slide.placeholders[1].text_frame
    tf.text = "• ScholarAI significantly reduces technical review time from hours to minutes."
    tf.add_paragraph().text = "• Multi-agentic workflows prove superior for deep document reasoning."
    tf.add_paragraph().text = "• Multimodal grounding drastically improves figure interpretability."
    for p in tf.paragraphs: p.font.color.rgb = RGBColor(255, 255, 255)

    # 12. Future Scope
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_background(slide, (15, 17, 23))
    slide.shapes.title.text = "Future Scope & Q&A"
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(76, 175, 130)
    tf = slide.placeholders[1].text_frame
    tf.text = "• Persistent Workspaces: SQLite-based browser session storage."
    tf.add_paragraph().text = "• Citation Generation: Native IEEE, BibTeX formats."
    tf.add_paragraph().text = "• Multi-Domain Tuning: Law, Medicine, and Law domain models."
    for p in tf.paragraphs: p.font.color.rgb = RGBColor(255, 255, 255)
    add_image_placeholder(slide, 4.0, 2)

    prs.save("ScholarAI_Final_Presentation.pptx")
    print("Presentation created successfully!")

if __name__ == "__main__":
    create_presentation()
